"""
Document Service — handles document upload, parsing, and chunking.
Supports PDF, PPTX, and DOCX formats.
Phase 1 provides the skeleton; full pipeline in Phase 4.
"""

import os
import logging
import uuid
from pathlib import Path
from typing import Optional, List, Dict, Tuple
from config import settings

logger = logging.getLogger("syllabus_ai")

# Supported file extensions
SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".ppt", ".docx", ".doc", ".txt"}

# Chunking settings
CHUNK_SIZE = 500  # words
CHUNK_OVERLAP = 100  # words


# Valid sections
VALID_SECTIONS = {"A", "B", "C"}


def get_document_dir(subject_code: str, section: str = None) -> Path:
    """Get the document storage directory for a subject and optional section.
    
    If section is provided (A, B, or C), returns documents/{subject}/{section}/
    Otherwise returns documents/{subject}/ for backward compatibility.
    """
    doc_dir = Path(settings.document_storage_dir) / subject_code.upper()
    if section and section.upper() in VALID_SECTIONS:
        doc_dir = doc_dir / section.upper()
    doc_dir.mkdir(parents=True, exist_ok=True)
    return doc_dir


def validate_file(filename: str) -> bool:
    """Check if the file extension is supported."""
    ext = Path(filename).suffix.lower()
    return ext in SUPPORTED_EXTENSIONS


async def extract_text(filepath: str) -> str:
    """
    Extract text from a document file.

    Supports: PDF, PPTX, DOCX, TXT
    """
    ext = Path(filepath).suffix.lower()

    try:
        if ext == ".pdf":
            return _extract_pdf(filepath)
        elif ext in (".pptx", ".ppt"):
            return _extract_pptx(filepath)
        elif ext in (".docx", ".doc"):
            return _extract_docx(filepath)
        elif ext == ".txt":
            return _extract_txt(filepath)
        else:
            raise ValueError(f"Unsupported file format: {ext}")
    except Exception as e:
        logger.error(f"Text extraction failed for {filepath}: {e}")
        raise


def _extract_pdf(filepath: str) -> str:
    """Extract text from PDF using PyMuPDF.
    Falls back to Gemini Vision OCR for scanned/handwritten pages.
    """
    import fitz  # PyMuPDF

    # ── Standard digital text extraction (existing behavior) ──
    text_parts = []
    with fitz.open(filepath) as doc:
        total_pages = len(doc)
        for page in doc:
            text_parts.append(page.get_text())

    full_text = "\n\n".join(text_parts)
    total_words = len(full_text.split())
    avg_words_per_page = total_words / max(total_pages, 1)

    # Sufficient text extracted → return as-is (unchanged behavior)
    if avg_words_per_page >= 20:
        return full_text

    # ── Low text density → likely scanned or handwritten ─────
    # Attempt Gemini Vision OCR if API key is configured
    if not settings.gemini_api_key:
        logger.warning(
            f"Low text ({avg_words_per_page:.0f} avg words/page) in "
            f"{Path(filepath).name}. Set GEMINI_API_KEY in .env for "
            f"handwritten/scanned note support."
        )
        return full_text

    logger.info(
        f"Low text density ({avg_words_per_page:.0f} words/page) in "
        f"{Path(filepath).name} — activating Gemini Vision OCR"
    )
    try:
        from services.ocr_service import extract_pdf_with_ocr
        return extract_pdf_with_ocr(filepath)
    except Exception as e:
        logger.warning(f"OCR fallback failed ({e}). Using standard extraction.")
        return full_text


def _extract_pptx(filepath: str) -> str:
    """Extract text from PowerPoint files."""
    from pptx import Presentation
    prs = Presentation(filepath)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_parts.append(shape.text)
    return "\n\n".join(text_parts)


def _extract_docx(filepath: str) -> str:
    """Extract text from Word documents."""
    from docx import Document
    doc = Document(filepath)
    text_parts = [para.text for para in doc.paragraphs if para.text.strip()]
    return "\n\n".join(text_parts)


def _extract_txt(filepath: str) -> str:
    """Read plain text files."""
    with open(filepath, "r", encoding="utf-8") as f:
        return f.read()


def chunk_text(
    text: str,
    chunk_size: int = CHUNK_SIZE,
    overlap: int = CHUNK_OVERLAP
) -> List[str]:
    """
    Split text into overlapping chunks by word count.

    Args:
        text: Full document text
        chunk_size: Words per chunk
        overlap: Overlapping words between chunks

    Returns:
        List of text chunks
    """
    words = text.split()
    chunks = []
    start = 0

    while start < len(words):
        end = start + chunk_size
        chunk = " ".join(words[start:end])
        if chunk.strip():
            chunks.append(chunk)
        start += chunk_size - overlap

    logger.info(f"Created {len(chunks)} chunks ({chunk_size}w, {overlap}w overlap)")
    return chunks


def generate_chunk_ids(
    subject_code: str,
    unit_number: int,
    filename: str,
    chunk_count: int
) -> List[str]:
    """Generate unique IDs for document chunks."""
    base = f"{subject_code}_u{unit_number}_{Path(filename).stem}"
    return [f"{base}_chunk_{i}" for i in range(chunk_count)]


async def list_documents(subject_code: str, section: str = None) -> List[Dict]:
    """List all uploaded documents for a subject (and optional section)."""
    doc_dir = get_document_dir(subject_code, section)
    documents = []

    for filepath in doc_dir.iterdir():
        if filepath.is_file() and filepath.suffix.lower() in SUPPORTED_EXTENSIONS:
            stat = filepath.stat()
            documents.append({
                "filename": filepath.name,
                "subject_code": subject_code,
                "section": section.upper() if section else None,
                "size_bytes": stat.st_size,
                "modified": stat.st_mtime
            })

    return documents
